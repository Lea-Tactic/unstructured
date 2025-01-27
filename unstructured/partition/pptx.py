from __future__ import annotations

import io
from tempfile import SpooledTemporaryFile
from typing import IO, Any, Iterator, List, Optional, Sequence, Tuple, Union

import pptx
from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.shapetree import _BaseGroupShapes  # pyright: ignore [reportPrivateUsage]
from pptx.slide import Slide
from pptx.text.text import _Paragraph  # pyright: ignore [reportPrivateUsage]

from unstructured.chunking.title import add_chunking_strategy
from unstructured.documents.elements import (
    Element,
    ElementMetadata,
    EmailAddress,
    ListItem,
    NarrativeText,
    PageBreak,
    Table,
    Text,
    Title,
    process_metadata,
)
from unstructured.file_utils.filetype import FileType, add_metadata_with_filetype
from unstructured.partition.common import (
    convert_ms_office_table_to_text,
    exactly_one,
    get_last_modified_date,
    get_last_modified_date_from_file,
)
from unstructured.partition.text_type import (
    is_email_address,
    is_possible_narrative_text,
    is_possible_title,
)
from unstructured.utils import lazyproperty


@process_metadata()
@add_metadata_with_filetype(FileType.PPTX)
@add_chunking_strategy()
def partition_pptx(
    filename: Optional[str] = None,
    file: Optional[IO[bytes]] = None,
    include_page_breaks: bool = True,
    metadata_filename: Optional[str] = None,
    include_metadata: bool = True,
    metadata_last_modified: Optional[str] = None,
    include_slide_notes: bool = False,
    chunking_strategy: Optional[str] = None,
    **kwargs: Any,
) -> List[Element]:
    """Partition PowerPoint document in .pptx format into its document elements.

    Parameters
    ----------
    filename
        A string defining the target filename path.
    file
        A file-like object using "rb" mode --> open(filename, "rb").
    include_page_breaks
        If True, includes a PageBreak element between slides
    metadata_filename
        The filename to use for the metadata. Relevant because partition_ppt() converts its
        (legacy) .ppt document to .pptx before partition. We want the filename of the original
        .ppt source file in the metadata.
    metadata_last_modified
        The last modified date for the document.
    include_slide_notes
        If True, includes the slide notes as element
    """
    # -- verify only one source-file argument was provided --
    exactly_one(filename=filename, file=file)

    # -- In Python <3.11 SpooledTemporaryFile does not implement ".seekable" which triggers an
    # -- exception when Zipfile tries to open it. Both the docx and pptx formats are zip archives,
    # -- so we need to work around that bug here.
    if isinstance(file, SpooledTemporaryFile):
        file.seek(0)
        file = io.BytesIO(file.read())

    source_file = file or filename
    assert source_file is not None

    return list(
        _PptxPartitioner.iter_presentation_elements(
            source_file,
            include_page_breaks,
            include_slide_notes,
            metadata_filename,
            metadata_last_modified,
        ),
    )


class _PptxPartitioner:  # pyright: ignore[reportUnusedClass]
    """Provides `.partition()` for PowerPoint 2007+ (.pptx) files."""

    def __init__(
        self,
        file: Union[str, IO[bytes]],
        # -- having default values for these arguments is not necessary for production uses because
        # -- this object is always created by the classmethod. However it simplifies constructing
        # -- this object in tests and makes them less sensitive to signature changes.
        include_page_breaks: bool = True,
        include_slide_notes: bool = False,
        metadata_filename: Optional[str] = None,
        metadata_last_modified: Optional[str] = None,
    ) -> None:
        self._file = file
        self._include_page_breaks = include_page_breaks
        self._include_slide_notes = include_slide_notes
        self._metadata_filename = metadata_filename
        self._metadata_last_modified = metadata_last_modified
        self._page_counter = 0

    @classmethod
    def iter_presentation_elements(
        cls,
        file: Union[str, IO[bytes]],
        include_page_breaks: bool,
        include_slide_notes: bool,
        metadata_filename: Optional[str],
        metadata_last_modified: Optional[str],
    ) -> Iterator[Element]:
        """Partition MS Word documents (.docx format) into its document elements."""
        return cls(
            file,
            include_page_breaks,
            include_slide_notes,
            metadata_filename,
            metadata_last_modified,
        )._iter_presentation_elements()

    def _iter_presentation_elements(self) -> Iterator[Element]:
        """Generate each document-element in presentation in document order."""
        # -- This implementation composes a collection of iterators into a "combined" iterator
        # -- return value using `yield from`. You can think of the return value as an Element
        # -- stream and each `yield from` as "add elements found by this function to the stream".
        # -- This is functionally analogous to declaring `elements: List[Element] = []` at the top
        # -- and using `elements.extend()` for the results of each of the function calls, but is
        # -- more perfomant, uses less memory (avoids producing and then garbage-collecting all
        # -- those small lists), is more flexible for later iterator operations like filter,
        # -- chain, map, etc. and is perhaps more elegant and simpler to read once you have the
        # -- concept of what it's doing. You can see the same pattern repeating in the "sub"
        # -- functions like `._iter_paragraph_elements()` where the "just return when done"
        # -- characteristic of a generator avoids repeated code to form interim results into lists.

        for slide in self._presentation.slides:
            yield from self._increment_page_number()
            yield from self._iter_maybe_slide_notes(slide)

            for shape in self._order_shapes(slide):
                if shape.has_table:
                    assert isinstance(shape, GraphicFrame)
                    yield from self._iter_table_element(shape)
                elif shape.has_text_frame:
                    assert isinstance(shape, Shape)
                    yield from self._iter_paragraph_elements(shape)
                # -- otherwise ditch it, this would include pictures, charts, connectors (lines),
                # -- and free-form shapes (squiggly lines). Lines don't have text.

    @lazyproperty
    def _filename(self) -> Optional[str]:
        """Suitable for use as metadata.filename, does not necessarily name source-file."""
        return (
            self._metadata_filename
            if self._metadata_filename
            else self._file
            if isinstance(self._file, str)
            else None
        )

    def _increment_page_number(self) -> Iterator[PageBreak]:
        """Increment page-number by 1 and generate a PageBreak element if enabled."""
        self._page_counter += 1
        # -- no page-break before first page --
        if self._page_counter < 2:
            return
        # -- only emit page-breaks when enabled --
        if self._include_page_breaks:
            yield PageBreak("")

    def _is_bulleted_paragraph(self, paragraph: _Paragraph) -> bool:
        """True when `paragraph` has a bullet-charcter prefix.

        Bullet characters in the openxml schema are represented by buChar.
        """
        # -- True when XPath returns a non-empty list (nodeset) --
        return bool(paragraph._p.xpath("./a:pPr/a:buChar"))

    def _iter_maybe_slide_notes(self, slide: Slide) -> Iterator[NarrativeText]:
        """Generate zero-or-one NarrativeText element for the slide-notes."""
        # -- only emit slide-notes elements when enabled --
        if not self._include_slide_notes:
            return

        # -- not all slides have a notes slide --
        if not slide.has_notes_slide:
            return

        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame

        # -- not all notes slides have a text-frame (it's created on first use) --
        if not notes_text_frame:
            return
        notes_text = notes_text_frame.text.strip()

        # -- not all notes text-frams contain text (if it's all deleted the text-frame remains) --
        if not notes_text:
            return

        yield NarrativeText(text=notes_text, metadata=self._text_metadata)

    def _iter_paragraph_elements(self, shape: Shape) -> Iterator[Element]:
        """Generate Text or subtype element for each paragraph in `shape`."""
        # NOTE(robinson) - avoid processing shapes that are not on the actual slide
        # NOTE - skip check if no top or left position (shape displayed top left)
        if (shape.top and shape.left) and (shape.top < 0 or shape.left < 0):
            return

        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text
            if text.strip() == "":
                continue
            if self._is_bulleted_paragraph(paragraph):
                yield ListItem(text=text, metadata=self._text_metadata)
            elif is_email_address(text):
                yield EmailAddress(text=text)
            elif is_possible_narrative_text(text):
                yield NarrativeText(text=text, metadata=self._text_metadata)
            elif is_possible_title(text):
                yield Title(text=text, metadata=self._text_metadata)
            else:
                yield Text(text=text, metadata=self._text_metadata)

    def _iter_table_element(self, graphfrm: GraphicFrame) -> Iterator[Table]:
        """Generate zero-or-one Table element for the table in `shape`.

        An empty table does not produce an element.
        """
        text_table = convert_ms_office_table_to_text(graphfrm.table, as_html=False).strip()
        if not text_table:
            return
        html_table = convert_ms_office_table_to_text(graphfrm.table, as_html=True)
        yield Table(text=text_table, metadata=self._table_metadata(html_table))

    @lazyproperty
    def _last_modified(self) -> Optional[str]:
        """Last-modified date suitable for use in element metadata."""
        # -- if this file was converted from another format, any last-modified date for the file
        # -- will be today, so we get it from the conversion step in `._metadata_last_modified`.
        if self._metadata_last_modified:
            return self._metadata_last_modified

        file = self._file

        # -- if the file is on the filesystem, get its date from there --
        if isinstance(file, str):
            return None if file.startswith("/tmp") else get_last_modified_date(file)

        # -- otherwise try getting it from the file-like object; this can work if `file` comes from
        # -- `with open(abc.pptx, "rb") as file:`, but I can't see folks doing that much when they
        # -- can just send us "abc.pptx" instead.
        return get_last_modified_date_from_file(file)

    def _order_shapes(self, slide: Slide) -> Sequence[BaseShape]:
        """Orders the shapes on `slide` from top to bottom and left to right."""

        def iter_shapes(shapes: _BaseGroupShapes) -> Iterator[BaseShape]:
            for shape in shapes:
                if isinstance(shape, GroupShape):
                    yield from iter_shapes(shape.shapes)
                else:
                    yield shape

        def sort_key(shape: BaseShape) -> Tuple[int, int]:
            return shape.top or 0, shape.left or 0

        return sorted(iter_shapes(slide.shapes), key=sort_key)

    @property
    def _page_number(self) -> Optional[int]:
        """The current page (slide) number."""
        return self._page_counter

    @lazyproperty
    def _presentation(self) -> Presentation:
        """The python-pptx `Presentation` object loaded from the provided source file."""
        return pptx.Presentation(self._file)

    def _table_metadata(self, text_as_html: str):
        """ElementMetadata instance suitable for use with Table element."""
        return ElementMetadata(
            filename=self._filename,
            last_modified=self._last_modified,
            page_number=self._page_number,
            text_as_html=text_as_html,
        )

    @property
    def _text_metadata(self):
        """ElementMetadata instance suitable for use with Text and subtypes."""
        return ElementMetadata(
            filename=self._filename,
            last_modified=self._last_modified,
            page_number=self._page_number,
        )
